# import libraries
import os
import json
from dotenv import load_dotenv
from azure.identity import DefaultAzureCredential
from azure.ai.documentintelligence import DocumentIntelligenceClient
from azure.ai.documentintelligence.models import AnalyzeResult
from azure.ai.documentintelligence.models import DocumentContentFormat

# load environment variables from .env file
load_dotenv()

# set endpoint from environment variable
endpoint = os.getenv("AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT")
credential = DefaultAzureCredential()

print("Endpoint:", endpoint)


def extract_polygon_and_images():
    """
    PPTX 파일에서 figure의 polygon 정보를 추출하고 JSON 파일로 저장합니다.
    또한 PPTX에 포함된 이미지를 추출합니다.
    """
    # PPTX 파일 경로
    file_path = "sample/sample-pptx.pptx"
    
    # 파일 존재 여부 확인
    if not os.path.exists(file_path):
        print(f"❌ Error: File not found at {file_path}")
        return
    
    print(f"📂 Processing file: {file_path}")
    print("=" * 80)

    document_intelligence_client = DocumentIntelligenceClient(
        endpoint=endpoint, credential=credential
    )

    # 로컬 파일을 바이너리로 읽어서 전송
    with open(file_path, "rb") as f:
        file_content = f.read()
        
    # PPTX 파일 분석 시작
    print("🔍 Analyzing document...")
    poller = document_intelligence_client.begin_analyze_document(
        "prebuilt-layout",
        body=file_content,
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        output_content_format=DocumentContentFormat.MARKDOWN
    )

    result: AnalyzeResult = poller.result()
    print("✅ Analysis complete!")
    print("=" * 80)

    # 출력 디렉토리 생성
    output_dir = "output"
    polygon_dir = os.path.join(output_dir, "pptx_polygons")
    images_dir = os.path.join(output_dir, "pptx_images")
    os.makedirs(polygon_dir, exist_ok=True)
    os.makedirs(images_dir, exist_ok=True)

    if not result.figures:
        print("❌ No figures found in the document.")
        return

    print(f"\n📊 Found {len(result.figures)} figures")
    print("=" * 80)

    # 전체 polygon 데이터를 담을 리스트
    all_polygons = []
    
    # python-pptx로 이미지 추출을 위해 PPTX 열기
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        pptx_available = True
    except ImportError:
        print("⚠️  python-pptx not installed. Image extraction will be skipped.")
        print("   To enable image extraction, run: pip install python-pptx")
        pptx_available = False

    # 각 figure 처리
    for idx, figure in enumerate(result.figures):
        if not figure.bounding_regions:
            continue
        
        figure_num = idx + 1
        bounding_region = figure.bounding_regions[0]
        page_num = bounding_region.page_number  # PPTX에서는 slide 번호
        polygon = bounding_region.polygon
        
        # polygon 정보 구조화
        polygon_data = {
            "figure_id": figure_num,
            "slide_number": page_num,
            "polygon_coordinates": polygon,
            "polygon_points": [
                {"x": polygon[i], "y": polygon[i+1]} 
                for i in range(0, len(polygon), 2)
            ],
            "caption": figure.caption.content if hasattr(figure, 'caption') and figure.caption else None,
            "elements_count": len(figure.elements) if hasattr(figure, 'elements') and figure.elements else 0
        }
        
        # bounding box 계산
        if len(polygon) >= 8:
            coords = [(polygon[i], polygon[i+1]) for i in range(0, len(polygon), 2)]
            x_coords = [c[0] for c in coords]
            y_coords = [c[1] for c in coords]
            
            polygon_data["bounding_box"] = {
                "x_min": min(x_coords),
                "y_min": min(y_coords),
                "x_max": max(x_coords),
                "y_max": max(y_coords),
                "width": max(x_coords) - min(x_coords),
                "height": max(y_coords) - min(y_coords)
            }
        
        all_polygons.append(polygon_data)
        
        # 개별 polygon JSON 파일로 저장
        figure_json_file = os.path.join(polygon_dir, f"figure_{figure_num:03d}_slide{page_num:02d}_polygon.json")
        with open(figure_json_file, "w", encoding="utf-8") as f:
            json.dump(polygon_data, f, indent=2, ensure_ascii=False)
        
        print(f"\n📐 Figure {figure_num} (Slide {page_num}):")
        print(f"   Polygon file: {figure_json_file}")
        
        if polygon_data.get("caption"):
            print(f"   Caption: {polygon_data['caption'][:60]}...")
    
    # PPTX에서 이미지 추출 (python-pptx 사용)
    if pptx_available:
        print("\n" + "=" * 80)
        print("🖼️  Extracting images from PPTX...")
        print("-" * 80)
        
        image_count = 0
        
        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_image_count = 0
            
            for shape_idx, shape in enumerate(slide.shapes):
                # Picture shape인 경우
                if hasattr(shape, "image"):
                    image_count += 1
                    slide_image_count += 1
                    
                    # 이미지 데이터 가져오기
                    image = shape.image
                    image_bytes = image.blob
                    ext = image.ext
                    
                    # 파일명 생성
                    filename = f"slide{slide_idx:02d}_image{slide_image_count:02d}.{ext}"
                    filepath = os.path.join(images_dir, filename)
                    
                    # 이미지 저장
                    with open(filepath, "wb") as f:
                        f.write(image_bytes)
                    
                    print(f"  ✅ Slide {slide_idx}, Image {slide_image_count}: {filename} ({len(image_bytes):,} bytes)")
        
        print("-" * 80)
        print(f"✅ Total images extracted from PPTX: {image_count}")
    
    # 전체 polygon 데이터를 하나의 JSON 파일로 저장
    all_polygons_file = os.path.join(output_dir, "pptx_all_polygons.json")
    summary_data = {
        "source_file": file_path,
        "total_figures": len(all_polygons),
        "total_slides": len(result.pages) if result.pages else 0,
        "figures": all_polygons
    }
    
    with open(all_polygons_file, "w", encoding="utf-8") as f:
        json.dump(summary_data, f, indent=2, ensure_ascii=False)
    
    print("\n" + "=" * 80)
    print(f"✅ Polygon extraction complete!")
    print(f"   Total figures processed: {len(all_polygons)}")
    print(f"   Individual polygon files: {polygon_dir}/")
    print(f"   All polygons summary: {all_polygons_file}")
    if pptx_available:
        print(f"   Extracted images: {images_dir}/")
    print("=" * 80)
    
    print("\n💡 Note:")
    print("   - Polygon coordinates are from Azure Document Intelligence")
    print("   - Images are extracted directly from PPTX file (all embedded images)")
    print("   - Polygon-based image cropping is not available for PPTX format")


if __name__ == "__main__":
    extract_polygon_and_images()
